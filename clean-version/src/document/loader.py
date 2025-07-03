"""
Document loading and processing functionality.
"""
import logging
from typing import List, Dict, Any, Optional, Union
from pathlib import Path
import tempfile
import os
from datetime import datetime

# Document processing imports
try:
    import pypdf
    from pypdf import PdfReader
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

try:
    import fitz  # PyMuPDF
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

import pandas as pd

from ..utils.exceptions import DocumentProcessingError
from ..utils.helpers import clean_text, safe_filename, get_file_extension, generate_hash
from config.settings import settings

logger = logging.getLogger(__name__)


class DocumentLoader:
    """Document loader for various file formats."""
    
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._load_pdf,
            '.docx': self._load_docx,
            '.pptx': self._load_pptx,
            '.xlsx': self._load_xlsx,
            '.xls': self._load_xlsx,
            '.txt': self._load_text,
            '.md': self._load_text,
            '.csv': self._load_csv,
            '.json': self._load_json
        }
        
    def load_document(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Load a document from file path.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dictionary containing document content and metadata
        """
        try:
            file_path = Path(file_path)
            
            if not file_path.exists():
                raise DocumentProcessingError(f"File not found: {file_path}")
            
            extension = get_file_extension(file_path.name)
            
            if extension not in self.supported_formats:
                raise DocumentProcessingError(f"Unsupported file format: {extension}")
            
            logger.info(f"Loading document: {file_path.name}")
            
            # Load document content
            content = self.supported_formats[extension](file_path)
            
            # Create document metadata
            metadata = {
                'filename': file_path.name,
                'filepath': str(file_path),
                'file_size': file_path.stat().st_size,
                'file_extension': extension,
                'loaded_at': datetime.now().isoformat(),
                'content_length': len(content),
                'document_hash': generate_hash(content)
            }
            
            return {
                'content': content,
                'metadata': metadata
            }
            
        except Exception as e:
            logger.error(f"Error loading document {file_path}: {str(e)}")
            raise DocumentProcessingError(f"Failed to load document: {str(e)}")
    
    def load_documents(self, file_paths: List[Union[str, Path]]) -> List[Dict[str, Any]]:
        """Load multiple documents."""
        documents = []
        errors = []
        
        for file_path in file_paths:
            try:
                document = self.load_document(file_path)
                documents.append(document)
            except Exception as e:
                error_msg = f"Failed to load {file_path}: {str(e)}"
                logger.error(error_msg)
                errors.append(error_msg)
        
        if errors:
            logger.warning(f"Failed to load {len(errors)} documents: {errors}")
        
        logger.info(f"Successfully loaded {len(documents)} documents")
        return documents
    
    def _load_pdf(self, file_path: Path) -> str:
        """Load PDF document."""
        if not HAS_PYPDF and not HAS_FITZ:
            raise DocumentProcessingError("PDF support requires pypdf or PyMuPDF")
        
        try:
            if HAS_FITZ:
                # Use PyMuPDF (better text extraction)
                doc = fitz.open(str(file_path))
                text = ""
                for page in doc:
                    text += page.get_text()
                doc.close()
                return clean_text(text)
            else:
                # Use pypdf
                with open(file_path, 'rb') as file:
                    pdf_reader = PdfReader(file)
                    text = ""
                    for page in pdf_reader.pages:
                        text += page.extract_text()
                return clean_text(text)
                
        except Exception as e:
            raise DocumentProcessingError(f"Failed to load PDF: {str(e)}")
    
    def _load_docx(self, file_path: Path) -> str:
        """Load DOCX document."""
        if not HAS_DOCX:
            raise DocumentProcessingError("DOCX support requires python-docx")
        
        try:
            doc = DocxDocument(str(file_path))
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return clean_text(text)
            
        except Exception as e:
            raise DocumentProcessingError(f"Failed to load DOCX: {str(e)}")
    
    def _load_pptx(self, file_path: Path) -> str:
        """Load PPTX document."""
        if not HAS_PPTX:
            raise DocumentProcessingError("PPTX support requires python-pptx")
        
        try:
            prs = Presentation(str(file_path))
            text = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return clean_text(text)
            
        except Exception as e:
            raise DocumentProcessingError(f"Failed to load PPTX: {str(e)}")
    
    def _load_xlsx(self, file_path: Path) -> str:
        """Load Excel document."""
        try:
            # Read all sheets
            df_dict = pd.read_excel(file_path, sheet_name=None)
            
            text = ""
            for sheet_name, df in df_dict.items():
                text += f"Sheet: {sheet_name}\n"
                text += df.to_string(index=False) + "\n\n"
            
            return clean_text(text)
            
        except Exception as e:
            raise DocumentProcessingError(f"Failed to load Excel: {str(e)}")
    
    def _load_text(self, file_path: Path) -> str:
        """Load text document."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            return clean_text(content)
            
        except UnicodeDecodeError:
            # Try different encodings
            encodings = ['latin-1', 'cp1252', 'iso-8859-1']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                    return clean_text(content)
                except UnicodeDecodeError:
                    continue
            raise DocumentProcessingError("Unable to decode text file")
            
        except Exception as e:
            raise DocumentProcessingError(f"Failed to load text: {str(e)}")
    
    def _load_csv(self, file_path: Path) -> str:
        """Load CSV document."""
        try:
            df = pd.read_csv(file_path)
            text = df.to_string(index=False)
            return clean_text(text)
            
        except Exception as e:
            raise DocumentProcessingError(f"Failed to load CSV: {str(e)}")
    
    def _load_json(self, file_path: Path) -> str:
        """Load JSON document."""
        try:
            import json
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
            
            # Convert JSON to readable text
            text = json.dumps(data, indent=2)
            return clean_text(text)
            
        except Exception as e:
            raise DocumentProcessingError(f"Failed to load JSON: {str(e)}")
    
    def load_from_bytes(self, content: bytes, filename: str) -> Dict[str, Any]:
        """Load document from bytes content."""
        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(suffix=get_file_extension(filename), delete=False) as tmp_file:
                tmp_file.write(content)
                tmp_file_path = tmp_file.name
            
            try:
                # Load document
                document = self.load_document(tmp_file_path)
                
                # Update metadata
                document['metadata']['filename'] = filename
                document['metadata']['source'] = 'uploaded'
                document['metadata']['original_size'] = len(content)
                
                return document
                
            finally:
                # Clean up temporary file
                os.unlink(tmp_file_path)
                
        except Exception as e:
            logger.error(f"Error loading document from bytes: {str(e)}")
            raise DocumentProcessingError(f"Failed to load document from bytes: {str(e)}")
    
    def get_supported_formats(self) -> List[str]:
        """Get list of supported file formats."""
        return list(self.supported_formats.keys())
    
    def is_supported_format(self, filename: str) -> bool:
        """Check if file format is supported."""
        extension = get_file_extension(filename)
        return extension in self.supported_formats
    
    def validate_document(self, document: Dict[str, Any]) -> bool:
        """Validate document structure."""
        required_keys = ['content', 'metadata']
        
        if not all(key in document for key in required_keys):
            return False
        
        if not isinstance(document['content'], str):
            return False
        
        if not isinstance(document['metadata'], dict):
            return False
        
        return True


class DocumentProcessor:
    """Document processor for handling multiple documents."""
    
    def __init__(self):
        self.loader = DocumentLoader()
        self.processed_documents = []
        
    def process_documents(self, file_paths: List[Union[str, Path]]) -> List[Dict[str, Any]]:
        """Process multiple documents."""
        try:
            logger.info(f"Processing {len(file_paths)} documents...")
            
            # Load documents
            documents = self.loader.load_documents(file_paths)
            
            # Process each document
            processed_docs = []
            for doc in documents:
                if self.loader.validate_document(doc):
                    processed_doc = self._process_single_document(doc)
                    processed_docs.append(processed_doc)
                else:
                    logger.warning(f"Invalid document structure: {doc.get('metadata', {}).get('filename', 'Unknown')}")
            
            self.processed_documents.extend(processed_docs)
            
            logger.info(f"Successfully processed {len(processed_docs)} documents")
            return processed_docs
            
        except Exception as e:
            logger.error(f"Error processing documents: {str(e)}")
            raise DocumentProcessingError(f"Failed to process documents: {str(e)}")
    
    def _process_single_document(self, document: Dict[str, Any]) -> Dict[str, Any]:
        """Process a single document."""
        try:
            # Add processing metadata
            document['metadata']['processed_at'] = datetime.now().isoformat()
            document['metadata']['processing_version'] = '1.0'
            
            # Add unique document ID
            document['id'] = generate_hash(
                document['content'] + 
                document['metadata']['filename'] + 
                str(document['metadata']['file_size'])
            )
            
            # Add document statistics
            content = document['content']
            document['metadata']['word_count'] = len(content.split())
            document['metadata']['character_count'] = len(content)
            document['metadata']['paragraph_count'] = len(content.split('\n\n'))
            
            return document
            
        except Exception as e:
            logger.error(f"Error processing single document: {str(e)}")
            raise DocumentProcessingError(f"Failed to process document: {str(e)}")
    
    def process_uploaded_files(self, uploaded_files: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Process uploaded files (for Streamlit integration)."""
        try:
            processed_docs = []
            
            for file_info in uploaded_files:
                filename = file_info['name']
                content = file_info['content']
                
                if self.loader.is_supported_format(filename):
                    document = self.loader.load_from_bytes(content, filename)
                    processed_doc = self._process_single_document(document)
                    processed_docs.append(processed_doc)
                else:
                    logger.warning(f"Unsupported file format: {filename}")
            
            self.processed_documents.extend(processed_docs)
            
            logger.info(f"Processed {len(processed_docs)} uploaded files")
            return processed_docs
            
        except Exception as e:
            logger.error(f"Error processing uploaded files: {str(e)}")
            raise DocumentProcessingError(f"Failed to process uploaded files: {str(e)}")
    
    def get_processed_documents(self) -> List[Dict[str, Any]]:
        """Get all processed documents."""
        return self.processed_documents
    
    def clear_processed_documents(self):
        """Clear all processed documents."""
        self.processed_documents = []
        logger.info("Cleared all processed documents")
    
    def get_processing_stats(self) -> Dict[str, Any]:
        """Get processing statistics."""
        if not self.processed_documents:
            return {
                'total_documents': 0,
                'total_content_length': 0,
                'supported_formats': self.loader.get_supported_formats()
            }
        
        total_content_length = sum(len(doc['content']) for doc in self.processed_documents)
        format_counts = {}
        
        for doc in self.processed_documents:
            ext = doc['metadata']['file_extension']
            format_counts[ext] = format_counts.get(ext, 0) + 1
        
        return {
            'total_documents': len(self.processed_documents),
            'total_content_length': total_content_length,
            'format_counts': format_counts,
            'supported_formats': self.loader.get_supported_formats(),
            'average_document_size': total_content_length / len(self.processed_documents)
        }
